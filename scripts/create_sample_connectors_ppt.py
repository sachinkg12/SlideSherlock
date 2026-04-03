#!/usr/bin/env python3
"""
Create a sample PPTX with shapes and connectors for testing G_native edge extraction.

The SlideSherlock pipeline only detects PowerPoint Connector shapes (not plain Line/Arrow).
This script uses python-pptx to add rectangles and connectors between them so that
ppt/slide_*.json has connectors and graphs/native/slide_*.json has non-empty edges.

Usage (from repo root):
  PYTHONPATH=. python scripts/create_sample_connectors_ppt.py [--output path.pptx]

Output: sample_connectors.pptx (or path given) in current directory.
"""

from __future__ import annotations

import argparse
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
except ImportError as e:
    print("python-pptx is required. pip install python-pptx", file=sys.stderr)
    raise SystemExit(1) from e


def create_sample_pptx(output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Blank slide (no title/body placeholders so we control all shapes)
    blank = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(blank)
    shapes = slide.shapes

    # Shape A: left rectangle
    left = Inches(0.5)
    top = Inches(2)
    w = Inches(1.5)
    h = Inches(0.8)
    rect_a = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    rect_a.text = "A"

    # Shape B: center rectangle
    left_b = Inches(4)
    top_b = Inches(2)
    rect_b = shapes.add_shape(MSO_SHAPE.RECTANGLE, left_b, top_b, w, h)
    rect_b.text = "B"

    # Shape C: right rectangle
    left_c = Inches(7.5)
    top_c = Inches(2)
    rect_c = shapes.add_shape(MSO_SHAPE.RECTANGLE, left_c, top_c, w, h)
    rect_c.text = "C"

    # Use integer EMU so saved file has int attributes (avoids python-pptx load error on float)
    def emu(v):
        if hasattr(v, "emu"):
            return Emu(int(v.emu))
        return Emu(int(v))

    # Connector A -> B: from right edge of A to left edge of B (same vertical center)
    center_y = emu(top + h / 2)
    begin_x = emu(left + w)
    end_x = emu(left_b)
    conn_ab = shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        begin_x, center_y,
        end_x, center_y,
    )

    # Connector B -> C
    begin_x_bc = emu(left_b + w)
    end_x_bc = emu(left_c)
    conn_bc = shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        begin_x_bc, center_y,
        end_x_bc, center_y,
    )

    # Optional: elbow connector from A to C (bottom path) to get a second connector type
    elbow_begin_x = emu(left + w / 2)
    elbow_begin_y = emu(top + h)
    elbow_end_x = emu(left_c + w / 2)
    elbow_end_y = emu(top_c + h)
    shapes.add_connector(
        MSO_CONNECTOR_TYPE.ELBOW,
        elbow_begin_x, elbow_begin_y,
        elbow_end_x, elbow_end_y,
    )

    os.makedirs(os.path.dirname(os.path.abspath(output_path)) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print("  Shapes: 3 rectangles (A, B, C)")
    print("  Connectors: 2 straight (A->B, B->C), 1 elbow (A->C bottom)")


def main() -> None:
    parser = argparse.ArgumentParser(description="Create sample PPTX with connectors for SlideSherlock.")
    parser.add_argument(
        "--output", "-o",
        default="sample_connectors.pptx",
        help="Output .pptx path (default: sample_connectors.pptx)",
    )
    args = parser.parse_args()
    if not args.output.endswith(".pptx"):
        args.output += ".pptx"
    create_sample_pptx(args.output)


if __name__ == "__main__":
    main()
