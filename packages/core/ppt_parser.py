"""
PPT parser: extract slide text, speaker notes, shapes (bbox, type, text runs, z_order),
connectors (endpoints, style, label), and groups from a .pptx file.
Output: list of per-slide dicts suitable for slide_i.json and DB metadata.
"""
from __future__ import annotations

import os
from typing import Any, Dict, List

try:
    from pptx import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.shapes.connector import Connector
    from pptx.shapes.group import GroupShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None  # type: ignore
    Connector = None  # type: ignore
    GroupShape = None  # type: ignore
    MSO_SHAPE_TYPE = None  # type: ignore


# EMU to something JSON-serializable (keep int EMU for precision)
def _emu_to_int(val: Any) -> int:
    if val is None:
        return 0
    if hasattr(val, "emu"):
        return int(val.emu)
    return int(val)


def _get_shape_bbox(shape: BaseShape) -> Dict[str, int]:
    """Return bbox as left, top, width, height in EMU."""
    return {
        "left": _emu_to_int(shape.left),
        "top": _emu_to_int(shape.top),
        "width": _emu_to_int(shape.width),
        "height": _emu_to_int(shape.height),
    }


def _get_text_runs(shape: BaseShape) -> List[Dict[str, str]]:
    """Extract text runs from a shape's text frame."""
    runs: List[Dict[str, str]] = []
    if not getattr(shape, "has_text_frame", True):
        return runs
    try:
        tf = shape.text_frame
        if tf is None:
            return runs
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text:
                    runs.append({"text": run.text.strip()})
    except Exception:
        pass
    return runs


def _shape_to_dict(shape: BaseShape, z_order: int, parent_prefix: str = "") -> Dict[str, Any]:
    """Serialize a single shape (non-connector, non-group) for JSON."""
    ppt_shape_id = f"{parent_prefix}{shape.shape_id}"
    out: Dict[str, Any] = {
        "ppt_shape_id": ppt_shape_id,
        "bbox": _get_shape_bbox(shape),
        "type": getattr(shape.shape_type, "name", str(shape.shape_type)),
        "text_runs": _get_text_runs(shape),
        "z_order": z_order,
    }
    return out


def _connector_to_dict(shape: Any, z_order: int, parent_prefix: str = "") -> Dict[str, Any]:
    """Serialize a connector shape: endpoints, style, label."""
    ppt_connector_id = f"{parent_prefix}conn_{shape.shape_id}"
    # Endpoints in EMU (relative to slide or group)
    try:
        begin_x = _emu_to_int(shape.begin_x) if hasattr(shape, "begin_x") else 0
        begin_y = _emu_to_int(shape.begin_y) if hasattr(shape, "begin_y") else 0
        end_x = _emu_to_int(shape.end_x) if hasattr(shape, "end_x") else 0
        end_y = _emu_to_int(shape.end_y) if hasattr(shape, "end_y") else 0
    except Exception:
        begin_x = begin_y = end_x = end_y = 0
    endpoints = {
        "begin": {"x": begin_x, "y": begin_y},
        "end": {"x": end_x, "y": end_y},
    }
    # Anchors: python-pptx doesn't always expose connection point indices; include if available
    anchors: Dict[str, Any] = {}
    if hasattr(shape, "begin_connected") and shape.begin_connected:
        anchors["begin_anchor"] = {"connected": True}
    if hasattr(shape, "end_connected") and shape.end_connected:
        anchors["end_anchor"] = {"connected": True}
    if anchors:
        endpoints["anchors"] = anchors

    style = ""
    try:
        if hasattr(shape, "connector_type"):
            style = getattr(shape.connector_type, "name", str(shape.connector_type))
    except Exception:
        pass

    label_runs = _get_text_runs(shape)
    label = " ".join(r["text"] for r in label_runs).strip() if label_runs else ""

    return {
        "ppt_connector_id": ppt_connector_id,
        "bbox": _get_shape_bbox(shape),
        "endpoints": endpoints,
        "style": style,
        "label": label,
        "z_order": z_order,
    }


def _group_to_dict(shape: GroupShape, z_order: int, parent_prefix: str = "") -> Dict[str, Any]:
    """Serialize a group: container with nested shapes (and connectors)."""
    ppt_shape_id = f"{parent_prefix}group_{shape.shape_id}"
    bbox = _get_shape_bbox(shape)
    children: List[Dict[str, Any]] = []
    prefix = f"{ppt_shape_id}_"
    for i, child in enumerate(shape.shapes):
        children.append(_extract_shape(child, i, prefix))
    return {
        "ppt_shape_id": ppt_shape_id,
        "bbox": bbox,
        "type": "GROUP",
        "text_runs": _get_text_runs(shape),
        "z_order": z_order,
        "children": children,
    }


def _extract_shape(shape: BaseShape, z_order: int, parent_prefix: str = "") -> Dict[str, Any]:
    """Dispatch to shape, connector, or group serializer."""
    if GroupShape is not None and isinstance(shape, GroupShape):
        return _group_to_dict(shape, z_order, parent_prefix)
    if Connector is not None and isinstance(shape, Connector):
        return _connector_to_dict(shape, z_order, parent_prefix)
    # Connector by type (some versions expose connector_type on shape)
    line_type = getattr(MSO_SHAPE_TYPE, "LINE", None) if MSO_SHAPE_TYPE else None
    if line_type and getattr(shape, "shape_type", None) == line_type and hasattr(shape, "begin_x"):
        return _connector_to_dict(shape, z_order, parent_prefix)
    return _shape_to_dict(shape, z_order, parent_prefix)


def _extract_notes(slide: Any) -> str:
    """Extract speaker notes text from a slide."""
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return ""
        tf = getattr(notes_slide, "notes_text_frame", None) or getattr(
            notes_slide, "text_frame", None
        )
        if tf is None:
            return ""
        parts = []
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text:
                    parts.append(run.text)
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _extract_slide_text(
    shapes: List[Dict[str, Any]], connectors: List[Dict[str, Any]], groups: List[Dict[str, Any]]
) -> str:
    """Concatenate all text from shapes, connectors (label), and groups for slide-level text."""
    parts = []

    def collect_text(obj: Dict[str, Any]) -> None:
        for r in obj.get("text_runs", []):
            t = r.get("text", "").strip()
            if t:
                parts.append(t)
        label = obj.get("label", "").strip()
        if label:
            parts.append(label)
        for c in obj.get("children", []):
            collect_text(c)

    for s in shapes + connectors + groups:
        collect_text(s)
    return " ".join(parts).strip()


def parse_pptx(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Parse a .pptx file and return a list of per-slide dicts.
    Each dict has: slide_index (1-based), slide_text, notes, shapes, connectors, groups.
    Shapes and connectors are in a unified list with a 'kind' discriminator; groups have children.
    """
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. pip install python-pptx")

    if not os.path.isfile(pptx_path):
        raise FileNotFoundError(pptx_path)

    prs = Presentation(pptx_path)
    slides_out: List[Dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        notes = _extract_notes(slide)
        shapes_out: List[Dict[str, Any]] = []
        connectors_out: List[Dict[str, Any]] = []
        groups_out: List[Dict[str, Any]] = []

        for z_order, shape in enumerate(slide.shapes):
            try:
                obj = _extract_shape(shape, z_order)
                if obj.get("ppt_connector_id"):
                    connectors_out.append(obj)
                elif obj.get("type") == "GROUP" or "children" in obj:
                    groups_out.append(obj)
                else:
                    shapes_out.append(obj)
            except Exception as e:
                # Log but continue; append a minimal placeholder so we don't lose z_order
                shapes_out.append(
                    {
                        "ppt_shape_id": f"err_{getattr(shape, 'shape_id', z_order)}",
                        "bbox": _get_shape_bbox(shape) if hasattr(shape, "left") else {},
                        "type": "UNKNOWN",
                        "text_runs": [],
                        "z_order": z_order,
                        "error": str(e),
                    }
                )

        # Slide-level text: from all shapes, connectors (label), and groups
        slide_text = _extract_slide_text(shapes_out, connectors_out, groups_out)

        slide_payload = {
            "slide_index": slide_idx,
            "slide_text": slide_text,
            "notes": notes,
            "shapes": shapes_out,
            "connectors": connectors_out,
            "groups": groups_out,
        }
        slides_out.append(slide_payload)

    return slides_out
