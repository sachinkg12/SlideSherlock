"""
Extract embedded images from PPTX and register as artifacts.
Outputs: jobs/{job_id}/images/slide_{i:03}/img_{k:02}.{ext}, jobs/{job_id}/images/index.json
"""
from __future__ import annotations

import hashlib
import json
from typing import Any, Dict, List

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    PRESENTATION_AVAILABLE = True
except ImportError:
    Presentation = None  # type: ignore
    MSO_SHAPE_TYPE = None  # type: ignore
    Emu = None  # type: ignore
    PRESENTATION_AVAILABLE = False

# EMU per inch (Office default)
EMU_PER_INCH = 914400

# Common slide dimensions in EMU (default 10" x 7.5")
DEFAULT_SLIDE_WIDTH_EMU = 914400
DEFAULT_SLIDE_HEIGHT_EMU = 685800


def _emu_to_float(val: Any) -> float:
    if val is None:
        return 0.0
    if hasattr(val, "emu"):
        return float(val.emu)
    return float(val)


def _stable_image_id(job_id: str, slide_index: int, ppt_shape_id: str) -> str:
    """Stable image_id: same ppt => same id across reruns."""
    payload = f"{job_id}|{slide_index}|IMAGE_ASSET|{ppt_shape_id}"
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def _ext_from_content_type(content_type: str) -> str:
    m = (content_type or "").lower()
    if "png" in m:
        return "png"
    if "jpeg" in m or "jpg" in m:
        return "jpg"
    if "gif" in m:
        return "gif"
    if "bmp" in m:
        return "bmp"
    if "tiff" in m or "tif" in m:
        return "tiff"
    return "png"


def _collect_picture_shapes(slide: Any) -> List[tuple]:
    """Return list of (shape, z_order, ppt_shape_id) for picture shapes, including from groups."""

    def recurse(shapes_iter: Any, z: int, prefix: str) -> List[tuple]:
        out: List[tuple] = []
        for i, shape in enumerate(shapes_iter):
            try:
                if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    ppt_shape_id = f"{prefix}{shape.shape_id}"
                    out.append((shape, z + i, ppt_shape_id))
                elif hasattr(shape, "shapes"):
                    gid = f"{prefix}group_{shape.shape_id}_"
                    for j, child in enumerate(shape.shapes):
                        if (
                            hasattr(child, "shape_type")
                            and child.shape_type == MSO_SHAPE_TYPE.PICTURE
                        ):
                            ppt_shape_id = f"{gid}{child.shape_id}"
                            out.append((child, z + i, ppt_shape_id))
            except Exception:
                continue
        return out

    return recurse(slide.shapes, 0, "")


def extract_images_from_pptx(
    pptx_path: str,
    job_id: str,
    minio_client: Any,
) -> Dict[str, Any]:
    """
    Extract embedded images from PPTX. For each picture shape:
    - Save bytes to jobs/{job_id}/images/slide_{i:03}/img_{k:02}.{ext}
    - Build index entry: image_id, ppt_shape_id, slide_index, bbox, normalized_bbox, z_index, mime, sha256, size
    Returns images index dict. Empty if no python-pptx or no images.
    """
    if not PRESENTATION_AVAILABLE or not Presentation:
        return {"images": [], "schema_version": "1.0", "job_id": job_id}

    if not pptx_path or not minio_client:
        return {"images": [], "schema_version": "1.0", "job_id": job_id}

    prs = Presentation(pptx_path)
    slide_width_emu = DEFAULT_SLIDE_WIDTH_EMU
    slide_height_emu = DEFAULT_SLIDE_HEIGHT_EMU
    try:
        slide_width_emu = int(prs.slide_width.emu)
        slide_height_emu = int(prs.slide_height.emu)
    except Exception:
        pass

    images: List[Dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        pictures = _collect_picture_shapes(slide)
        for k, (shape, z_order, ppt_shape_id) in enumerate(pictures):
            try:
                image = shape.image
                blob = image.blob
                if not blob:
                    continue
                content_type = getattr(image, "content_type", None) or "image/png"
                ext = getattr(image, "ext", None) or _ext_from_content_type(content_type)
                sha256_hash = hashlib.sha256(blob).hexdigest()
                size = len(blob)

                image_id = _stable_image_id(job_id, slide_idx, ppt_shape_id)

                left_emu = _emu_to_float(getattr(shape, "left", None) or 0)
                top_emu = _emu_to_float(getattr(shape, "top", None) or 0)
                width_emu = _emu_to_float(getattr(shape, "width", None) or 0)
                height_emu = _emu_to_float(getattr(shape, "height", None) or 0)

                bbox = {
                    "x": left_emu,
                    "y": top_emu,
                    "w": width_emu,
                    "h": height_emu,
                }

                norm_x = left_emu / slide_width_emu if slide_width_emu else 0
                norm_y = top_emu / slide_height_emu if slide_height_emu else 0
                norm_w = width_emu / slide_width_emu if slide_width_emu else 0
                norm_h = height_emu / slide_height_emu if slide_height_emu else 0
                normalized_bbox = {"x": norm_x, "y": norm_y, "w": norm_w, "h": norm_h}

                storage_path = f"jobs/{job_id}/images/slide_{slide_idx:03d}/img_{k:02d}.{ext}"
                minio_client.put(storage_path, blob, content_type)

                entry = {
                    "image_id": image_id,
                    "ppt_shape_id": ppt_shape_id,
                    "slide_index": slide_idx,
                    "bbox": bbox,
                    "normalized_bbox": normalized_bbox,
                    "z_index": z_order,
                    "mime": content_type,
                    "ext": ext,
                    "sha256": sha256_hash,
                    "size": size,
                    "uri": storage_path,
                }
                images.append(entry)
            except Exception:
                continue

    index_payload = {
        "schema_version": "1.0",
        "job_id": job_id,
        "images": images,
    }
    index_path = f"jobs/{job_id}/images/index.json"
    minio_client.put(
        index_path,
        json.dumps(index_payload, indent=2).encode("utf-8"),
        "application/json",
    )
    return index_payload
