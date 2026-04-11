"""
Tests for ppt_parser.py: _emu_to_int, _get_shape_bbox, _get_text_runs,
_extract_slide_text, _extract_notes, parse_pptx (mocked and real pptx).
"""
from __future__ import annotations

import os
import sys
from unittest.mock import MagicMock, patch

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from ppt_parser import (
    _emu_to_int,
    _get_shape_bbox,
    _get_text_runs,
    _extract_slide_text,
    _extract_notes,
    _shape_to_dict,
    _connector_to_dict,
    _extract_shape,
    parse_pptx,
)


# ---------------------------------------------------------------------------
# _emu_to_int
# ---------------------------------------------------------------------------

def test_emu_to_int_plain_int():
    assert _emu_to_int(914400) == 914400


def test_emu_to_int_none_returns_zero():
    assert _emu_to_int(None) == 0


def test_emu_to_int_emu_attribute():
    obj = MagicMock()
    obj.emu = 457200
    assert _emu_to_int(obj) == 457200


def test_emu_to_int_float_truncated():
    assert _emu_to_int(3.9) == 3


# ---------------------------------------------------------------------------
# _get_shape_bbox
# ---------------------------------------------------------------------------

def _make_shape(left=10, top=20, width=100, height=50):
    shape = MagicMock()
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height
    return shape


def test_get_shape_bbox_returns_dict():
    shape = _make_shape(10, 20, 100, 50)
    bbox = _get_shape_bbox(shape)
    assert bbox == {"left": 10, "top": 20, "width": 100, "height": 50}


# ---------------------------------------------------------------------------
# _get_text_runs
# ---------------------------------------------------------------------------

def _make_shape_with_text(texts):
    """Build a mock shape whose text_frame has paragraphs containing runs."""
    shape = MagicMock()
    shape.has_text_frame = True
    paras = []
    for t in texts:
        run = MagicMock()
        run.text = t
        para = MagicMock()
        para.runs = [run]
        paras.append(para)
    shape.text_frame.paragraphs = paras
    return shape


def test_get_text_runs_returns_stripped_texts():
    shape = _make_shape_with_text(["  Hello  ", "  World  "])
    runs = _get_text_runs(shape)
    assert runs == [{"text": "Hello"}, {"text": "World"}]


def test_get_text_runs_skips_empty_runs():
    # The source checks `if run.text:` before stripping; whitespace-only strings
    # are truthy so they pass the guard and become {"text": ""} after strip.
    # Only truly empty strings ("") are excluded by the guard.
    shape = _make_shape_with_text(["", "Keep"])
    runs = _get_text_runs(shape)
    # "" is falsy so skipped; "Keep" passes and stays
    assert runs == [{"text": "Keep"}]


def test_get_text_runs_no_text_frame():
    shape = MagicMock()
    shape.has_text_frame = False
    runs = _get_text_runs(shape)
    assert runs == []


# ---------------------------------------------------------------------------
# _extract_slide_text
# ---------------------------------------------------------------------------

def test_extract_slide_text_combines_all_sources():
    shapes = [{"text_runs": [{"text": "Title"}], "label": "", "children": []}]
    connectors = [{"text_runs": [], "label": "connects", "children": []}]
    groups = [
        {
            "text_runs": [{"text": "Group"}],
            "label": "",
            "children": [{"text_runs": [{"text": "Child"}], "label": "", "children": []}],
        }
    ]
    result = _extract_slide_text(shapes, connectors, groups)
    assert "Title" in result
    assert "connects" in result
    assert "Group" in result
    assert "Child" in result


def test_extract_slide_text_empty_inputs():
    assert _extract_slide_text([], [], []) == ""


# ---------------------------------------------------------------------------
# _extract_notes
# ---------------------------------------------------------------------------

def test_extract_notes_no_notes_slide():
    slide = MagicMock()
    slide.has_notes_slide = False
    assert _extract_notes(slide) == ""


def test_extract_notes_with_text():
    run = MagicMock()
    run.text = "Speaker note text"
    para = MagicMock()
    para.runs = [run]
    notes_slide = MagicMock()
    notes_slide.notes_text_frame.paragraphs = [para]
    slide = MagicMock()
    slide.has_notes_slide = True
    slide.notes_slide = notes_slide
    result = _extract_notes(slide)
    assert result == "Speaker note text"


def test_extract_notes_exception_returns_empty():
    slide = MagicMock()
    slide.has_notes_slide = True
    slide.notes_slide = None  # causes AttributeError accessing .notes_text_frame
    assert _extract_notes(slide) == ""


# ---------------------------------------------------------------------------
# parse_pptx – missing file / missing dependency
# ---------------------------------------------------------------------------

def test_parse_pptx_missing_file_raises():
    # When python-pptx is not installed, RuntimeError is raised before the file
    # check. When it is installed, FileNotFoundError is raised. Both are valid.
    with pytest.raises((FileNotFoundError, RuntimeError)):
        parse_pptx("/nonexistent/path/deck.pptx")


# ---------------------------------------------------------------------------
# parse_pptx – real .pptx via python-pptx (skipped if not installed)
# ---------------------------------------------------------------------------

def test_parse_pptx_returns_slide_list(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Hello"
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))

    result = parse_pptx(str(pptx_path))
    assert isinstance(result, list)
    assert len(result) == 1
    slide_dict = result[0]
    assert slide_dict["slide_index"] == 1
    assert "shapes" in slide_dict
    assert "connectors" in slide_dict
    assert "groups" in slide_dict


def test_parse_pptx_slide_text_populated(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "My Title"
    pptx_path = tmp_path / "titled.pptx"
    prs.save(str(pptx_path))

    result = parse_pptx(str(pptx_path))
    assert result[0]["slide_text"]  # non-empty string expected
