"""Shared fixtures for SlideSherlock core tests."""

from __future__ import annotations

import os
import sys
import tempfile
from pathlib import Path
from unittest.mock import MagicMock

import pytest

# Ensure packages/core is on the path for all tests
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))


@pytest.fixture
def tmp_dir(tmp_path):
    """Provide a temporary directory that auto-cleans."""
    return tmp_path


@pytest.fixture
def sample_slide_data():
    """Minimal slide data structure used across tests."""
    return {
        "slide_index": 0,
        "title": "Test Slide",
        "body_text": "This is a test slide body.",
        "notes": "Speaker notes here.",
        "shapes": [],
        "images": [],
    }


@pytest.fixture
def sample_script_segment():
    """Minimal script segment for pipeline tests."""
    return {
        "claim_id": "c1",
        "slide_index": 0,
        "text": "This slide shows the architecture overview.",
        "evidence_ids": ["e1"],
        "entity_ids": ["ent1"],
    }


@pytest.fixture
def sample_evidence_index():
    """Minimal evidence index mapping."""
    return {
        "e1": {
            "evidence_id": "e1",
            "slide_index": 0,
            "source": "body_text",
            "text": "Architecture overview of the system.",
        },
    }


@pytest.fixture
def sample_graph():
    """Minimal unified graph structure."""
    return {
        "nodes": {
            "ent1": {"id": "ent1", "label": "Architecture", "type": "concept"},
        },
        "edges": [],
    }


@pytest.fixture
def mock_llm_provider():
    """Mock LLM provider that returns canned responses."""
    provider = MagicMock()
    provider.complete.return_value = "Mocked LLM response."
    provider.complete_json.return_value = {"result": "ok"}
    return provider


@pytest.fixture
def mock_vision_provider():
    """Mock vision provider for image understanding tests."""
    provider = MagicMock()
    provider.describe_image.return_value = {
        "description": "A diagram showing system components.",
        "objects": ["box", "arrow", "label"],
        "text_detected": ["Component A", "Component B"],
    }
    return provider


@pytest.fixture
def sample_pptx_path(tmp_path):
    """Create a minimal .pptx file for parsing tests."""
    try:
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = "Test Presentation"
        body = slide.placeholders[1]
        body.text = "Test body content for slide 1."
        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return path
    except ImportError:
        pytest.skip("python-pptx not installed")
